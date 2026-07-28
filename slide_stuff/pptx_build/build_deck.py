"""Build output/coax-projections.pptx from the Svelte deck's measured layout.

Pipeline (see README in this directory):

    bun slides/scripts/dump-slide-data.ts   # numbers, via the TS coax model
    bun slides/scripts/dump-layout.mjs      # geometry, via headless Chrome
    uv run pptx_build/build_deck.py         # this script

Every displayed number comes from the same TypeScript physics library the deck
uses, so the pptx can never drift from the slides. Nothing here re-derives it.
"""

from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from slide_stuff.pptx_build.render import Scale, draw_slide

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "output"
LAYOUT = OUTPUT / "slide-layout.json"
DATA = OUTPUT / "slide-data.json"
DECK = OUTPUT / "coax-projections.pptx"

# 16:9 widescreen, the PowerPoint default.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

BLANK_LAYOUT = 6


def count_kinds(nodes: list[dict], tally: Counter | None = None) -> Counter:
    tally = Counter() if tally is None else tally
    for node in nodes:
        tally[node["kind"]] += 1
        if node["kind"] == "group":
            count_kinds(node["children"], tally)
    return tally


def build() -> Path:
    layout = json.loads(LAYOUT.read_text())
    data = json.loads(DATA.read_text())

    images = {label: OUTPUT / file for label, file in layout["xsecs"].items()}
    missing = [str(p) for p in images.values() if not p.exists()]
    if missing:
        raise FileNotFoundError(f"cross-section images not exported: {missing}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[BLANK_LAYOUT]

    for spec in layout["slides"]:
        slide = prs.slides.add_slide(blank)
        sc = Scale(sx=SLIDE_W / spec["width"], sy=SLIDE_H / spec["height"])
        draw_slide(slide, spec["nodes"], sc, images)

        tally = count_kinds(spec["nodes"])
        print(f"  slide_{spec['index']}: " + ", ".join(f"{v} {k}" for k, v in sorted(tally.items())))

    core = prs.core_properties
    core.title = data["deckTitle"]
    core.subject = data["title"]["kicker"]

    prs.save(DECK)
    return DECK


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
