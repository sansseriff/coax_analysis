"""Draw browser-measured geometry onto a python-pptx slide.

The deck's real layout engine is Chrome: `slides/scripts/dump-layout.mjs` walks
the rendered SvelteKit deck and records, in CSS px relative to the slide box,
every background/border box, every *rendered line* of text, and the cable
cross-section placeholders. This module is the other half of that bridge — it
knows nothing about coax, only how to turn those records into shapes.

Two conventions make the result faithful without any hand-tuning:

* Text is emitted one shape per rendered line, sized to the line's tight ink
  box and centred vertically in it. Chrome has already resolved line breaking,
  inline margins and flex/grid positioning, so wrapping is switched off and the
  glyphs land where the browser put them.
* A CSS box may carry four differently-styled borders, which a pptx autoshape
  cannot express. When at least three sides agree we draw that as the shape's
  outline and overlay only the odd side(s) as thin rectangles.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU_PER_PT = 12700

# Weights below this render as regular; pptx has no variable-weight axis.
BOLD_AT = 600

# Rounded corners eat into a one-sided border bar; pull the bar in by this
# fraction of the corner radius so it doesn't poke past the curve.
BAR_INSET = 0.45

ALIGN = {
    "right": PP_ALIGN.RIGHT,
    "end": PP_ALIGN.RIGHT,
    "center": PP_ALIGN.CENTER,
}


@dataclass(frozen=True)
class Scale:
    """CSS px -> EMU, one axis at a time (the frame is not exactly 16:9)."""

    sx: float
    sy: float

    def x(self, px: float) -> Emu:
        return Emu(int(round(px * self.sx)))

    def y(self, px: float) -> Emu:
        return Emu(int(round(px * self.sy)))

    def pt(self, px: float) -> Pt:
        return Pt(px * self.sx / EMU_PER_PT)


def _rgb(c: dict[str, float]) -> RGBColor:
    return RGBColor(int(c["r"]), int(c["g"]), int(c["b"]))


def _apply_alpha(elm, color: dict[str, float]) -> None:
    """Attach <a:alpha> to the srgbClr under `elm` (a solidFill parent)."""
    if color["a"] >= 1:
        return
    srgb = elm.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(round(color["a"] * 100000))))


def _fill(shape, color: dict[str, float] | None) -> None:
    if color is None:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    _apply_alpha(shape._element.spPr, color)


def _outline(shape, border: dict[str, Any] | None, sc: Scale) -> None:
    if border is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = _rgb(border["color"])
    shape.line.width = sc.x(border["w"])
    _apply_alpha(shape._element.spPr.find(qn("a:ln")), border["color"])


def _rect(shapes, x, y, w, h, sc: Scale, radius: float = 0.0):
    """A rectangle (rounded when `radius`), with fill and line left unset."""
    w = max(w, 0.05)
    h = max(h, 0.05)
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = shapes.add_shape(kind, sc.x(x), sc.y(y), sc.x(w), sc.y(h))
    if radius:
        shape.adjustments[0] = min(radius / min(w, h), 0.5)

    # A new autoshape references the theme's shape style, which carries a drop
    # shadow. Drop the reference and pin an empty effect list so no renderer
    # reintroduces one; every colour on this deck is set explicitly anyway.
    sp = shape._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)
    shape.shadow.inherit = False
    return shape


def _uniform(borders: dict[str, Any]) -> tuple[dict | None, list[str]]:
    """Split borders into a dominant outline and the odd sides left over.

    Returns `(outline, extra_sides)`. `outline` is only offered when three or
    four sides share a width and colour — otherwise every present side is drawn
    as its own bar, because a pptx outline is all-or-nothing.
    """
    present = {k: v for k, v in borders.items() if v}
    if not present:
        return None, []

    groups: dict[tuple, list[str]] = {}
    for side, b in present.items():
        key = (b["w"], tuple(b["color"].values()))
        groups.setdefault(key, []).append(side)

    best = max(groups.values(), key=len)
    if len(best) >= 3:
        return present[best[0]], [s for s in present if s not in best]
    return None, list(present)


def _border_bar(shapes, rect: dict, side: str, border: dict, radius: float, sc: Scale) -> None:
    x, y, w, h = rect["x"], rect["y"], rect["w"], rect["h"]
    bw = border["w"]
    inset = radius * BAR_INSET

    if side in ("top", "bottom"):
        bx, bw_, by, bh = x + inset, w - 2 * inset, y if side == "top" else y + h - bw, bw
    else:
        bx, bw_, by, bh = (x if side == "left" else x + w - bw), bw, y + inset, h - 2 * inset

    bar = _rect(shapes, bx, by, bw_, bh, sc)
    _fill(bar, border["color"])
    _outline(bar, None, sc)


def draw_box(shapes, node: dict, sc: Scale) -> None:
    rect, radius = node["rect"], node.get("radius", 0)
    outline, extras = _uniform(node["borders"])

    if node["bg"] is not None or outline is not None:
        shape = _rect(shapes, rect["x"], rect["y"], rect["w"], rect["h"], sc, radius)
        _fill(shape, node["bg"])
        _outline(shape, outline, sc)

    for side in extras:
        _border_bar(shapes, rect, side, node["borders"][side], radius, sc)


def draw_text(shapes, node: dict, sc: Scale) -> None:
    """One text box per rendered line; its runs re-flow within that line."""
    rect = node["rect"]
    box = shapes.add_textbox(sc.x(rect["x"]), sc.y(rect["y"]), sc.x(rect["w"]), sc.y(rect["h"]))

    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # The rect is the line's ink box, so centring reproduces CSS half-leading.
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    para = tf.paragraphs[0]
    para.alignment = ALIGN.get(node["align"], PP_ALIGN.LEFT)

    for spec in node["runs"]:
        run = para.add_run()
        run.text = spec["text"]
        font = run.font
        font.name = spec["font"]
        font.size = sc.pt(spec["sizePx"])
        font.bold = spec["weight"] >= BOLD_AT
        font.italic = spec["italic"]
        font.color.rgb = _rgb(spec["color"])

        if spec["letterSpacingPx"]:
            spacing = spec["letterSpacingPx"] * sc.sx / EMU_PER_PT
            run.font._rPr.set("spc", str(int(round(spacing * 100))))


def draw_picture(shapes, node: dict, image: Path, sc: Scale) -> None:
    rect = node["rect"]
    shapes.add_picture(
        str(image), sc.x(rect["x"]), sc.y(rect["y"]), sc.x(rect["w"]), sc.y(rect["h"])
    )


def draw_nodes(shapes, nodes: list[dict], sc: Scale, images: dict[str, Path]) -> None:
    """Paint nodes in document order, which is also pptx z-order.

    `shapes` is a slide's shape tree or a group's, so groups nest naturally.
    """
    for node in nodes:
        kind = node["kind"]
        if kind == "box":
            draw_box(shapes, node, sc)
        elif kind == "text":
            draw_text(shapes, node, sc)
        elif kind == "svg":
            draw_picture(shapes, node, images[node["label"]], sc)
        elif kind == "group":
            group = shapes.add_group_shape()
            group.name = node["name"]
            draw_nodes(group.shapes, node["children"], sc, images)
            # Extents are recalculated as each child is added, but a nested group
            # grows after its parent last measured it — so re-measure bottom-up.
            group._element.recalculate_extents()
        else:
            raise ValueError(f"unknown node kind: {kind!r}")


def draw_slide(slide, nodes: list[dict], sc: Scale, images: dict[str, Path]) -> None:
    draw_nodes(slide.shapes, nodes, sc, images)
