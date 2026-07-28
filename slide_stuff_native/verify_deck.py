"""Structural checks on the built deck.

    uv run verify_deck.py

Asserts the deliverable's hard requirements: 16:9, seven slides, a real
GraphicFrame table on slide 6, grouped cross-sections, native text, nothing
off-canvas, and every figure in the JSON actually present on its slide.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

ROOT = Path(__file__).parent
DATA = json.loads((ROOT / "data" / "slide-data.json").read_text(encoding="utf-8"))
prs = Presentation(ROOT / "output" / "coax-deck.pptx")

fails: list[str] = []


def check(ok: bool, msg: str):
    if not ok:
        fails.append(msg)


def texts(slide) -> str:
    out = []

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(sh.shapes)
            elif sh.has_text_frame:
                out.append(sh.text_frame.text)
            elif sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        out.append(c.text)

    walk(slide.shapes)
    return "\n".join(out)


def groups(slide) -> int:
    return sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.GROUP)


# ── Deck-level ──
check(prs.slide_width == Emu(12192000), f"slide width {prs.slide_width} != 16:9 13.333in")
check(prs.slide_height == Emu(6858000), f"slide height {prs.slide_height}")
check(len(prs.slides) == 7, f"{len(prs.slides)} slides, expected 7")

W, H = prs.slide_width, prs.slide_height
for i, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        if sh.left < -9525 or sh.top < -9525 or sh.left + sh.width > W + 9525 or sh.top + sh.height > H + 9525:
            fails.append(f"slide {i}: '{sh.shape_type}' out of bounds "
                         f"({sh.left/914400:.2f},{sh.top/914400:.2f} "
                         f"{sh.width/914400:.2f}×{sh.height/914400:.2f})")

# ── Slide 0: title ──
t0 = texts(prs.slides[0])
check("Insertion Loss & Heat-Load Projections" in t0, "slide 0 missing headline")
check("50 Ω stranded micro-coax · cryostat wiring" in t0, "slide 0 missing eyebrow")
for c in DATA["cables"]:
    check(c["name"] in t0, f"slide 0 missing {c['name']}")
    check(c["description"] in t0, f"slide 0 missing description of {c['name']}")
    check(c["z0"] in t0, f"slide 0 missing Z0 {c['z0']}")
    check(c["odMm"] in t0 and c["odMils"] in t0, f"slide 0 missing OD of {c['name']}")
    check(c["alpha2"] in t0, f"slide 0 missing alpha {c['alpha2']}")
for L in DATA["meta"]["lengthsLabel"].split(", "):
    check(L in t0, f"slide 0 missing length {L}")

# ── Slides 1–5: one per length ──
for i, L in enumerate(DATA["lengths"], start=1):
    slide = prs.slides[i]
    t = texts(slide)
    check(L["title"] in t, f"slide {i} missing title {L['title']!r}")
    for cid in ("samtec", "n12"):
        d = L["cables"][cid]
        c = next(x for x in DATA["cables"] if x["id"] == cid)
        check(c["name"] in t, f"slide {i} missing {c['name']}")
        for field in ("loss2", "loss4"):
            check(d[field] in t, f"slide {i} missing {cid} {field}={d[field]}")
        for field in ("q40", "q41"):
            check(d[field]["v"] in t, f"slide {i} missing {cid} {field}={d[field]['v']}")
            check(d[field]["u"] in t, f"slide {i} missing unit {d[field]['u']}")
    tr = DATA["tradeoff"]
    check(f"{tr['lossRatio']}× lower loss" in t, f"slide {i} missing loss ratio")
    check(f"{tr['heatRatio']}× lower heat" in t, f"slide {i} missing heat ratio")
    # Two cross-sections + spec pills + figure rows, all grouped.
    check(groups(slide) >= 14, f"slide {i} has {groups(slide)} groups, expected >= 14")

# ── Slide 6: a real table ──
s6 = prs.slides[6]
frames = [sh for sh in s6.shapes if sh.has_table]
check(len(frames) == 1, f"slide 6 has {len(frames)} tables, expected exactly 1")
if frames:
    tbl = frames[0].table
    check(frames[0].shape_type == MSO_SHAPE_TYPE.TABLE, "slide 6 table is not a GraphicFrame table")
    check(len(tbl.rows) == 8, f"table has {len(tbl.rows)} rows, expected 8 (3 header + 5 lengths)")
    check(len(tbl.columns) == 11, f"table has {len(tbl.columns)} cols, expected 11")
    merged = sum(1 for r in tbl.rows for c in r.cells if c.is_merge_origin)
    check(merged >= 5, f"table has {merged} merge origins, expected >= 5 (group headers)")

    body = "\n".join(c.text for r in tbl.rows for c in r.cells)
    for c in DATA["cables"]:
        check(c["name"] in body, f"table missing group header {c['name']}")
    for row in DATA["table"]:
        check(row["length"] in body, f"table missing row {row['length']}")
        for cell in row["cols"]:
            check(cell["text"].replace(" ", "") in body.replace(" ", ""),
                  f"table missing cell {cell['text']!r}")

t6 = texts(s6)
check(DATA["meta"]["harnessLabel"] in t6, "slide 6 does not label the harness")

# ── Report ──
if fails:
    print(f"FAIL — {len(fails)} problem(s):")
    for f in fails:
        print("  •", f)
    sys.exit(1)
print(f"PASS — 7 slides, 16:9, real table (8×11), "
      f"{sum(groups(s) for s in prs.slides)} grouped shapes, all figures present")
