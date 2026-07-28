"""Build the 7-slide cryogenic micro-coax deck with python-pptx.

    uv run build_deck.py

Reads data/slide-data.json (produced by `bun run ts/dump.ts`) and writes
output/coax-deck.pptx. Every displayed string is taken verbatim from that JSON.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from crosssection import common_scale, cross_section
from table import summary_table
from theme import (
    ACCENT,
    BG,
    HAIRLINE,
    INK,
    INK_DIM,
    INK_FAINT,
    MARGIN,
    MONO,
    PANEL,
    SLIDE_H,
    SLIDE_W,
    eyebrow,
    figure_row,
    footnote,
    kicker,
    para,
    pill,
    rect,
    rule,
    run,
    set_background,
    soft_shadow,
    textbox,
)

ROOT = Path(__file__).parent
DATA = json.loads((ROOT / "data" / "slide-data.json").read_text(encoding="utf-8"))
OUT = ROOT / "output" / "coax-deck.pptx"

CABLE_IDS = ["samtec", "n12"]
CABLES = {c["id"]: c for c in DATA["cables"]}

# Two equal columns, used identically on the title slide and the length slides.
COL_W = (SLIDE_W - 2 * MARGIN - Inches(0.38)) / 2
COL_X = [MARGIN, MARGIN + COL_W + Inches(0.38)]
PAD = Inches(0.34)


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    return slide


def slide_head(sh, brow: str, title: str, right: str | None = None):
    """Eyebrow, title, and the hairline that closes the header band."""
    eyebrow(sh, MARGIN, Inches(0.44), Inches(8), brow)
    _, tf = textbox(sh, MARGIN, Inches(0.72), Inches(9.5), Inches(0.55))
    run(para(tf, True), title, 26, color=INK, bold=True, spacing=-0.5)
    if right:
        _, tf = textbox(sh, SLIDE_W - MARGIN - Inches(5.4), Inches(0.90), Inches(5.4), Inches(0.3))
        run(para(tf, True, align=PP_ALIGN.RIGHT), right, 8.2, color=INK_FAINT, spacing=1.1)
    rule(sh, MARGIN, Inches(1.26), SLIDE_W - 2 * MARGIN)


# ── Slide 0 — title ──────────────────────────────────────────────────────────


def slide_title(prs):
    s = blank(prs)
    sh = s.shapes

    eyebrow(sh, MARGIN, Inches(1.02), Inches(9), "50 Ω stranded micro-coax · cryostat wiring", size=11)

    _, tf = textbox(sh, MARGIN, Inches(1.36), Inches(12), Inches(0.85))
    run(para(tf, True, line=1.0), "Insertion Loss & Heat-Load Projections", 38, color=INK,
        bold=True, spacing=-0.9)

    # Lede: the cables carry their accent, the five lengths carry the emphasis.
    _, tf = textbox(sh, MARGIN, Inches(2.36), Inches(11.6), Inches(0.8))
    p = para(tf, True, line=1.4)
    run(p, "Two 50 Ω candidates for a 208-line harness — ", 13, color=INK_DIM)
    run(p, CABLES["samtec"]["name"], 13, color=ACCENT["samtec"], bold=True)
    run(p, " and ", 13, color=INK_DIM)
    run(p, CABLES["n12"]["name"], 13, color=ACCENT["n12"], bold=True)
    run(p, " — projected across five assembly lengths: ", 13, color=INK_DIM)
    run(p, DATA["meta"]["lengthsLabel"] + ".", 13, color=INK, bold=True)

    # ── Two spec cards, each with an accent rail down its left edge ──
    card_y, card_h = Inches(3.42), Inches(2.36)
    for i, cid in enumerate(CABLE_IDS):
        c = CABLES[cid]
        accent = ACCENT[cid]
        x = COL_X[i]

        soft_shadow(rect(sh, x, card_y, COL_W, card_h, fill=PANEL, line=HAIRLINE, lw=0.75, radius=0.05))
        rect(sh, x, card_y, Inches(0.05), card_h, fill=accent)

        tx = x + PAD
        tw = COL_W - 2 * PAD
        _, tf = textbox(sh, tx, card_y + Inches(0.28), tw, Inches(0.32))
        run(para(tf, True), c["name"], 17, color=accent, bold=True, spacing=-0.2)

        _, tf = textbox(sh, tx, card_y + Inches(0.66), tw, Inches(0.4))
        run(para(tf, True, line=1.3), c["description"], 9.6, color=INK_DIM)

        # The three length-independent specs.
        ry = card_y + Inches(1.14)
        rh = Inches(0.34)
        big, sml = 13, 9.2
        specs = [
            ("Z₀", [(c["z0"], big, INK, True), (" Ω", sml, INK_FAINT, False)]),
            ("Cable OD", [(c["odMm"], big, INK, True), (" mm", sml, INK_FAINT, False),
                          ("  ·  ", sml, INK_FAINT, False),
                          (c["odMils"], big, INK, True), (" mils", sml, INK_FAINT, False)]),
            ("α @ 2 GHz", [(c["alpha2"], big, INK, True), (" dB/m", sml, INK_FAINT, False)]),
        ]
        for j, (label, vruns) in enumerate(specs):
            figure_row(sh, tx, ry + j * Inches(0.38), tw, rh, label,
                       label_color=INK_FAINT, underline=j < len(specs) - 1,
                       value_runs=vruns)

    rule(sh, MARGIN, Inches(6.16), SLIDE_W - 2 * MARGIN)
    footnote(
        sh,
        "Method — conductor loss after Pozar (Microwave Engineering, §2.7) plus dielectric loss, with a "
        "per-conductor stranding factor K_s calibrated to the TCF-3450F insertion-loss curve over 1–4 GHz. "
        "Heat load is steady-state conduction across each cryostat stage's ΔT, with lay-angle path corrections. "
        "Insertion loss is quoted at the 40 K → 4 K stage (T̄ ≈ 22 K); cross-sections are drawn to each preset's "
        "real geometry at a common scale.",
        y=Inches(6.36),
    )
    return s


# ── Slides 1–5 — one per assembly length ─────────────────────────────────────


def cable_column(sh, x, y, w, h, cid: str, L: dict, scale: float):
    c = CABLES[cid]
    accent = ACCENT[cid]
    d = L["cables"][cid]

    soft_shadow(rect(sh, x, y, w, h, fill=BG, line=HAIRLINE, lw=0.75, radius=0.035))
    rect(sh, x, y, w, Inches(0.05), fill=accent)

    tx = x + PAD
    tw = w - 2 * PAD

    # Cross-section, top-right of the panel.
    r_frame = Inches(0.58)
    cx = int(x + w - PAD - r_frame)
    cy = int(y + Inches(0.32) + r_frame)
    cross_section(sh, c["geom"], cx, cy, scale)

    _, tf = textbox(sh, tx, y + Inches(0.26), tw - Inches(1.5), Inches(0.32))
    run(para(tf, True), c["name"], 16, color=accent, bold=True, spacing=-0.2)

    _, tf = textbox(sh, tx, y + Inches(0.62), tw - Inches(1.5), Inches(0.5))
    run(para(tf, True, line=1.3), c["description"], 9.2, color=INK_DIM)

    # Length-independent specs, as subordinate chips.
    _, pw = pill(sh, tx, y + Inches(1.34), "Z₀", c["z0"] + " Ω")
    pill(sh, tx + pw + Inches(0.10), y + Inches(1.34), "OD",
         c["odMm"] + " mm · " + c["odMils"] + " mils")

    rule(sh, tx, y + Inches(1.86), tw)

    # The two figure groups this deck exists to compare.
    kicker(sh, tx, y + Inches(1.98), tw, "RF insertion loss")
    figure_row(sh, tx, y + Inches(2.20), tw, Inches(0.50), "2 GHz", d["loss2"], "dB", 30, INK)
    figure_row(sh, tx, y + Inches(2.80), tw, Inches(0.32), "4 GHz", d["loss4"], "dB", 16,
               INK_DIM, underline=False)

    kicker(sh, tx, y + Inches(3.32), tw, "Conducted heat load · per cable")
    figure_row(sh, tx, y + Inches(3.54), tw, Inches(0.50), "40 K → 4 K",
               d["q40"]["v"], d["q40"]["u"], 30, INK)
    figure_row(sh, tx, y + Inches(4.14), tw, Inches(0.32), "4 K → 1 K",
               d["q41"]["v"], d["q41"]["u"], 16, INK_DIM, underline=False)


def slide_length(prs, L: dict):
    s = blank(prs)
    sh = s.shapes

    slide_head(sh, "CABLE LENGTH", L["title"], "PER-CABLE FIGURES · 50 Ω · 40 K → 4 K STAGE")

    panel_y, panel_h = Inches(1.42), Inches(4.72)
    scale = common_scale([CABLES[i]["geom"] for i in CABLE_IDS], Inches(0.58))
    for i, cid in enumerate(CABLE_IDS):
        cable_column(sh, COL_X[i], panel_y, COL_W, panel_h, cid, L, scale)

    # ── Trade-off line: identical on all five slides, and that is the point. ──
    t = DATA["tradeoff"]
    _, tf = textbox(sh, MARGIN, Inches(6.30), SLIDE_W - 2 * MARGIN, Inches(0.3))
    p = para(tf, True, align=PP_ALIGN.CENTER)
    run(p, "Trade-off  ·  ", 11, color=INK_DIM, font=MONO)
    run(p, t["lossWinner"], 11, color=ACCENT["samtec"], bold=True, font=MONO)
    run(p, f"  {t['lossRatio']}× lower loss", 11, color=INK_DIM, font=MONO)
    run(p, "  ·  ", 11, color=INK_FAINT, font=MONO)
    run(p, t["heatWinner"], 11, color=ACCENT["n12"], bold=True, font=MONO)
    run(p, f"  {t['heatRatio']}× lower heat", 11, color=INK_DIM, font=MONO)

    rule(sh, MARGIN, Inches(6.72), SLIDE_W - 2 * MARGIN)
    footnote(
        sh,
        "Both cables are 50 Ω. Model — Pozar conductor loss + dielectric loss with a per-conductor stranding "
        "factor K_s calibrated to the TCF-3450F insertion-loss curve (1–4 GHz); heat load is steady-state "
        "conduction across each stage's ΔT with lay-angle path corrections. Insertion loss is quoted at the "
        "40 K → 4 K stage. The two ratios are length-independent — they are the same on every length slide.",
        y=Inches(6.88),
    )
    return s


# ── Slide 6 — summary ────────────────────────────────────────────────────────


def slide_summary(prs):
    s = blank(prs)
    sh = s.shapes

    slide_head(sh, "SUMMARY", "All five lengths at a glance",
               f"FULL HARNESS · {DATA['meta']['harnessLabel']}")

    summary_table(sh, DATA, CABLES, CABLE_IDS, y=Inches(1.68))

    rule(sh, MARGIN, Inches(6.44), SLIDE_W - 2 * MARGIN)
    footnote(
        sh,
        "Insertion loss quoted at 2 GHz, 4–40 K stage — see the per-length slides for the 4 GHz figures. "
        "Heat load is the conducted Q̇ across each stage's ΔT, shown per cable and for the full harness "
        f"({DATA['meta']['harnessLabel']}). The bundle columns are the per-cable figure × {DATA['meta']['harness']}; "
        "they are the ones that set the cryostat budget.",
        y=Inches(6.60),
    )
    return s


# ── Assemble ─────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    for L in DATA["lengths"]:
        slide_length(prs, L)
    slide_summary(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)} — {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
