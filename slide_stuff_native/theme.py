"""Palette, type scale and low-level shape primitives for the coax deck.

The look is ported from the HTML/Svelte version of this deck: white ground,
near-black text, one accent per cable, Inter for prose and JetBrains Mono for
every number. Colours below were sampled from those exported PNGs.

Everything here is presentation-only. No physical quantity is computed in
Python and no number is rounded here — display strings arrive pre-formatted
from ts/dump.ts.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── Canvas ───────────────────────────────────────────────────────────────────

# Exact 16:9 — Inches(13.333) lands 305 EMU short of PowerPoint's widescreen width.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Inches(0.62)


# ── Palette (sampled from the HTML export) ───────────────────────────────────


def rgb(h: str) -> RGBColor:
    return RGBColor.from_string(h.lstrip("#").upper())


def mix(a: str, b: str, t: float) -> str:
    """Blend two hex colours, t=0 → a, t=1 → b."""
    ai, bi = int(a.lstrip("#"), 16), int(b.lstrip("#"), 16)
    out = 0
    for shift in (16, 8, 0):
        ca, cb = (ai >> shift) & 0xFF, (bi >> shift) & 0xFF
        out |= round(ca + (cb - ca) * t) << shift
    return f"#{out:06X}"


BG = "#FFFFFF"
PANEL = "#F7F9FC"  # card fill, and the tint on the bundle columns
PANEL_HI = "#F1F4F8"
HAIRLINE = "#E2E6EC"

INK = "#161B24"  # figures and headlines
INK_DIM = "#4A5263"  # prose, row labels
INK_FAINT = "#8A93A6"  # units, captions, footnotes

# One accent per cable, used on every slide: card rule, name, table group header.
ACCENT = {
    "samtec": "#2A6FB0",
    "n12": "#B9701A",
}
BRAND = ACCENT["samtec"]  # eyebrows

# Cross-section materials. The Samtec core is silver-plated *copper*; the N12 is
# brass — they must not look like the same metal.
METAL_FILL = {"spc": "#B87333", "brass": "#C9A64A", "copper": "#B87333"}
METAL_LINE = {"spc": "#95591F", "brass": "#9C7C2C", "copper": "#95591F"}
DIELECTRIC = "#C9D6E9"
DIELECTRIC_LINE = "#C2D1E8"
JACKET = "#FFFFFF"
JACKET_LINE = "#E2E6ED"

FONT = "Inter"
MONO = "JetBrains Mono"
MONO_ADV = 0.6  # JetBrains Mono advance width, in ems — lets us size pills exactly


# ── Text primitives ──────────────────────────────────────────────────────────


def textbox(shapes, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    """A real, wrapping text frame with zero insets."""
    tb = shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf


def para(tf, first: bool = False, align=PP_ALIGN.LEFT, space_before: int = 0, line: float | None = None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    return p


def run(p, text, size, color=INK, bold=False, italic=False, spacing=0, font=FONT):
    """Add a styled run. `spacing` is letter-spacing in points."""
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = rgb(color)
    if spacing:
        r._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))
    return r


def smart_upper(s: str) -> str:
    """Uppercase ASCII only — 'α' must not become the Latin 'A'."""
    return "".join(c.upper() if c.isascii() else c for c in s)


def kicker(shapes, x, y, w, text, color=INK_FAINT, size=8.6, font=FONT):
    """Uppercase, letter-spaced section label."""
    _, tf = textbox(shapes, x, y, w, Pt(size * 1.8))
    run(para(tf, True), smart_upper(text), size, color=color, bold=True, spacing=1.2, font=font)
    return tf


def eyebrow(shapes, x, y, w, text, size=9.5):
    """The mono, accent-coloured line above every slide title."""
    _, tf = textbox(shapes, x, y, w, Pt(size * 1.9))
    run(para(tf, True), text, size, color=BRAND, bold=False, spacing=0.7, font=MONO)
    return tf


def mono_width(text: str, size: float) -> float:
    """Width of a mono string in points — exact, because the font is monospaced."""
    return len(text) * size * MONO_ADV


# ── Shape primitives ─────────────────────────────────────────────────────────


def _plain(shape):
    """Kill the theme drop shadow.

    `shadow.inherit = False` only adds an empty <a:effectLst/>, which PowerPoint
    honours but LibreOffice does not — LO keeps rendering the preset shadow off
    the shape's <p:style> effectRef. So we strip that effectRef outright. Every
    autoshape (rules, cross-section dots, pills, accent rails) goes through here
    and comes out shadowless; cards opt back in via soft_shadow().
    """
    shape.shadow.inherit = False
    style = shape._element.find(qn("p:style"))
    if style is not None:
        ref = style.find(qn("a:effectRef"))
        if ref is not None:
            ref.set("idx", "0")  # 0 = no effect style
            for child in list(ref):
                ref.remove(child)
    return shape


def soft_shadow(shape, blur=13, dist=4, alpha=13, color="1B2433"):
    """A wide, faint drop shadow for cards. blur/dist in points, alpha in %."""
    spPr = shape._element.spPr
    effectLst = spPr.get_or_add_effectLst()
    for child in list(effectLst):
        effectLst.remove(child)
    shdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur * 12700)),
        "dist": str(int(dist * 12700)),
        "dir": "5400000",  # straight down
        "rotWithShape": "0",
    })
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": color})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))}))
    shdw.append(clr)
    effectLst.append(shdw)
    return shape


def rect(shapes, x, y, w, h, fill=None, line=None, lw=0.75, radius=None):
    if radius is None:
        sh = shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    else:
        sh = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
        sh.adjustments[0] = radius
    _plain(sh)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = rgb(line)
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    sh.text_frame.word_wrap = True
    return sh


def circle(shapes, cx, cy, r, fill=None, line=None, lw=0.5):
    sh = shapes.add_shape(MSO_SHAPE.OVAL, int(cx - r), int(cy - r), int(2 * r), int(2 * r))
    _plain(sh)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = rgb(line)
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def rule(shapes, x, y, w, color=HAIRLINE, h=Pt(0.75)):
    return rect(shapes, x, y, w, int(h), fill=color)


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


# Glyphs built from a combining mark. Inter drops the mark when LibreOffice
# shapes it; JetBrains Mono composes it correctly. Letter-spacing also breaks
# the cluster, so these runs never carry `spc`.
COMBINING = ("Q̇", "T̄")


def _split_combining(text: str):
    """Split prose into (segment, font) pairs, isolating combining-mark glyphs."""
    parts = [(text, FONT)]
    for sym in COMBINING:
        out = []
        for seg, font in parts:
            if font is not FONT or sym not in seg:
                out.append((seg, font))
                continue
            chunks = seg.split(sym)
            for i, ch in enumerate(chunks):
                if i:
                    out.append((sym, MONO))
                if ch:
                    out.append((ch, FONT))
        parts = out
    return parts


def footnote(shapes, text, y, w=None):
    _, tf = textbox(shapes, MARGIN, y, w or (SLIDE_W - 2 * MARGIN), Inches(0.5))
    p = para(tf, True, line=1.35)
    for seg, font in _split_combining(text):
        run(p, seg, 7.8, color=INK_FAINT, font=font)
    return tf


# ── Composite elements ───────────────────────────────────────────────────────


def pill(shapes, x, y, label: str, value: str, size=8.6, h=Inches(0.29)):
    """A small rounded chip: faint mono label, then the value in ink.

    Returns (group, width) so the caller can lay several out in a row.
    """
    pad = Pt(10)
    # Mono, so the text width is exact; a little slack keeps it off the pill edge.
    w = int(Pt(mono_width(label + " " + value, size) + 2) + 2 * pad)

    g = shapes.add_group_shape()
    gs = g.shapes
    rect(gs, x, y, w, h, fill=PANEL, line=HAIRLINE, lw=0.75, radius=0.42)
    _, tf = textbox(gs, x + pad, y, w - 2 * pad, h, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    p = para(tf, True)
    run(p, label + " ", size, color=INK_FAINT, font=MONO)
    run(p, value, size, color=INK_DIM, bold=True, font=MONO)
    return g, w


def figure_row(
    shapes,
    x,
    y,
    w,
    h,
    label: str,
    value: str = "",
    unit: str = "",
    value_size: float = 13,
    value_color: str = INK,
    label_size: float = 9.6,
    label_color: str = INK_DIM,
    underline: bool = True,
    value_runs: list | None = None,
):
    """One line of a figure group: label on the left, big value + small unit right.

    Value and unit live in one text frame so they share a baseline, and the whole
    row is a GroupShape so label, number and unit travel together in PowerPoint.
    Pass `value_runs` as [(text, size, color, bold), …] for a compound value such
    as "0.56 mm · 22.1 mils".
    """
    g = shapes.add_group_shape()
    gs = g.shapes

    _, tf = textbox(gs, x, y, w * 0.45, h, anchor=MSO_ANCHOR.BOTTOM)
    run(para(tf, True), label, label_size, color=label_color, font=MONO)

    _, tf = textbox(gs, x, y, w, h, anchor=MSO_ANCHOR.BOTTOM, wrap=False)
    p = para(tf, True, align=PP_ALIGN.RIGHT, line=0.95)
    if value_runs:
        for text, size, color, bold in value_runs:
            run(p, text, size, color=color, bold=bold, spacing=-0.3 if bold else 0, font=MONO)
    else:
        run(p, value, value_size, color=value_color, bold=True, spacing=-0.3, font=MONO)
        if unit:
            # Unit always smaller than its number, but never so small it disappears.
            run(p, " " + unit, max(value_size * 0.42, 8.2), color=INK_FAINT, font=MONO)

    if underline:
        rule(gs, x, y + h + Pt(3), w)
    return g
